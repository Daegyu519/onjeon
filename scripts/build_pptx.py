#!/usr/bin/env python3
"""덱 HTML → **편집 가능한** .pptx (네이티브 텍스트 상자 · 도형 · 표).

    uv pip install -p .venv -e ".[docs]"
    .venv/bin/python scripts/build_pptx.py

`docs/deck/온전_기술설명서.html`이 여전히 유일한 원본이다. 이 스크립트는 그 HTML을
파싱해서 같은 문구를 pptx 도형으로 다시 그린다 — 문구를 파이썬으로 옮겨 적으면
HTML과 pptx가 갈라지고, 그러면 어느 쪽이 제출본인지 알 수 없게 된다.

이미지 슬라이드를 버린 이유: 심사자·팀원이 한 글자도 못 고친다. 대신 조판이
브라우저와 완전히 같지는 않다(줄바꿈 위치가 몇 군데 다를 수 있다). 픽셀이 정본인
쪽은 PDF가 맡고(`build_deck.py`, 같은 HTML을 헤드리스 캡처), pptx는 편집을 맡는다.

좌표계: HTML 덱이 1920×1080(vmin=10.8px) 기준이라 그대로 쓴다.
  px → 인치 = px/144,  px → 포인트 = px/2.
"""
from __future__ import annotations

import math
import re
import struct
import sys
from dataclasses import dataclass, field
from html.parser import HTMLParser
from pathlib import Path
from typing import Callable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
DECK = ROOT / "docs" / "deck" / "온전_기술설명서.html"
OUT = ROOT / "dist-docs" / "온전_기술설명서.pptx"

W, H = 1920.0, 1080.0
VMIN = H / 100
PAD_X, PAD_Y = 6.4 * VMIN, 5.2 * VMIN
CW = W - 2 * PAD_X          # 콘텐츠 폭 1782px
CH = H - 2 * PAD_Y          # 콘텐츠 높이 968px

# 조판 추정치는 실제 렌더보다 조금 짧게 나온다. 카드·표가 글자를 물지 않게 여유를 준다.
SLACK = 1.08

C = {
    "bg": "1B1917", "surface": "252119", "line": "3A352C",
    "text": "F6F3EC", "muted": "A99F8F", "dim": "7B7263",
    "kb": "FFBC00", "kb_deep": "BF7D00", "kb_ink": "26282B",
    "bad": "E5674A", "good": "7FB069",
    "on": "2E2718",  # .card.on 의 노란 틴트를 배경에 미리 섞은 값
}
MONO = "Consolas"


# ────────────────────────────────────────────────────────────── 단위/측정

def IN(px: float) -> Emu:
    return Inches(px / 144.0)


def PT(px: float):
    return Pt(px / 2.0)


def char_w(ch: str, size: float) -> float:
    o = ord(ch)
    if o > 0x1100:            # 한글·한자·전각 (실측 렌더가 1em보다 조금 좁다)
        return size * 0.95
    if ch == " ":
        return size * 0.27
    if ch in ".,·:;'!|()[]":
        return size * 0.30
    return size * 0.54


def text_w(s: str, size: float) -> float:
    return sum(char_w(c, size) for c in s)


Line = list[tuple[str, dict]]


def text_h(lines: list[Line], w: float, base: float, lh: float) -> float:
    total = 0.0
    for ln in lines:
        if not ln:
            total += base * lh
            continue
        sz = max(st.get("size", base) for _, st in ln)
        wide = sum(text_w(t, st.get("size", base)) for t, st in ln)
        total += max(1, math.ceil(wide / w - 1e-6)) * sz * lh
    return total * SLACK


# ────────────────────────────────────────────────────────────── 미니 DOM

@dataclass
class Node:
    tag: str
    attrs: dict = field(default_factory=dict)
    children: list = field(default_factory=list)

    @property
    def cls(self) -> set[str]:
        return set(self.attrs.get("class", "").split())

    def kids(self) -> list["Node"]:
        return [c for c in self.children if isinstance(c, Node)]

    def find(self, *classes: str) -> list["Node"]:
        want = set(classes)
        hit = []
        for k in self.kids():
            if want & k.cls or k.tag in want:
                hit.append(k)
            hit += k.find(*classes)
        return hit


VOID = {"br", "img", "meta", "link", "hr", "input"}


class _Parser(HTMLParser):
    def __init__(self):
        super().__init__(convert_charrefs=True)
        self.root = Node("#root")
        self.stack = [self.root]

    def handle_starttag(self, tag, attrs):
        n = Node(tag, dict(attrs))
        self.stack[-1].children.append(n)
        if tag not in VOID:
            self.stack.append(n)

    def handle_startendtag(self, tag, attrs):
        self.stack[-1].children.append(Node(tag, dict(attrs)))

    def handle_endtag(self, tag):
        if tag in VOID:
            return
        for i in range(len(self.stack) - 1, 0, -1):
            if self.stack[i].tag == tag:
                del self.stack[i:]
                return

    def handle_data(self, data):
        self.stack[-1].children.append(data)


# ────────────────────────────────────────────────────────── 인라인 스타일

_WS = re.compile(r"[ \t\n\r\f\v]+")       # \xa0(&nbsp;)는 남긴다


def css_color(v: str) -> str | None:
    v = v.strip()
    m = re.fullmatch(r"var\(--([\w-]+)\)", v)
    if m:
        return C.get(m.group(1).replace("-", "_"))
    if v.startswith("#"):
        return v[1:].upper().ljust(6, "0")[:6] if len(v) > 4 else "".join(c * 2 for c in v[1:])
    return None


def css_len(v: str) -> float | None:
    m = re.fullmatch(r"([\d.]+)vmin", v.strip())
    return float(m.group(1)) * VMIN if m else None


def inline_style(raw: str) -> dict:
    out: dict = {}
    for part in raw.split(";"):
        if ":" not in part:
            continue
        k, v = part.split(":", 1)
        k = k.strip()
        if k == "color" and (c := css_color(v)):
            out["color"] = c
        elif k == "font-size" and (s := css_len(v)) is not None:
            out["size"] = s
        elif k == "font-weight":
            out["bold"] = v.strip() in ("700", "800", "bold")
    return out


def style_of(n: Node) -> dict:
    s: dict = {}
    if n.tag in ("b", "strong"):
        s.update(bold=True, color=C["text"])
    if n.tag == "code":
        s.update(mono=True, color=C["kb"], rel=0.92)
    c = n.cls
    if "hl" in c:
        s["color"] = C["kb"]
    if "c" in c:                      # .formula .c — 보조 설명 줄
        s.update(color=C["dim"], size=1.6 * VMIN)
    if raw := n.attrs.get("style"):
        s.update(inline_style(raw))
    return s


def merge(base: dict, new: dict) -> dict:
    out = dict(base)
    new = dict(new)
    if (rel := new.pop("rel", None)) is not None:
        out["size"] = base.get("size", 20.0) * rel
    out.update(new)
    return out


def lines_of(n: Node, base: dict, skip: Callable[[Node], bool] | None = None) -> list[Line]:
    """노드의 인라인 내용을 <br> 기준 줄 목록으로. skip이 True인 자식은 건너뛴다."""
    out: list[Line] = [[]]

    def walk(node: Node, st: dict):
        for ch in node.children:
            if isinstance(ch, str):
                t = _WS.sub(" ", ch)
                if t.strip(" ") or (t == " " and out[-1]):
                    out[-1].append((t, st))
            elif ch.tag == "br":
                out.append([])
            elif skip and skip(ch):
                continue
            else:
                walk(ch, merge(st, style_of(ch)))

    walk(n, base)
    clean: list[Line] = []
    for ln in out:
        while ln and not ln[0][0].strip(" "):
            ln.pop(0)
        while ln and not ln[-1][0].strip(" "):
            ln.pop()
        if ln:
            ln[0] = (ln[0][0].lstrip(" "), ln[0][1])
            ln[-1] = (ln[-1][0].rstrip(" "), ln[-1][1])
        clean.append(ln)
    while len(clean) > 1 and not clean[-1]:
        clean.pop()
    return clean


def txt(n: Node) -> str:
    return "".join(
        _WS.sub(" ", c) if isinstance(c, str) else (" " if c.tag == "br" else txt(c))
        for c in n.children
    ).strip()


# ────────────────────────────────────────────────────────────── 그리기

def fill_tf(tf, lines: list[Line], base: float, color: str, lh: float,
            align=PP_ALIGN.LEFT, bold=False, mono=False, tracking: float = 0.0):
    tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = lh
        p.space_after = 0
        p.space_before = 0
        for t, st in ln or [("", {})]:
            r = p.add_run()
            r.text = t
            f = r.font
            f.size = PT(st.get("size", base))
            f.bold = st.get("bold", bold)
            f.color.rgb = RGBColor.from_string(st.get("color", color))
            if st.get("mono", mono):
                f.name = MONO
            if tracking:
                r.font._rPr.set("spc", str(int(tracking * st.get("size", base) / 2 * 100)))


def text_box(sl, x, y, w, h, lines, base, color, lh=1.4, **kw):
    tb = sl.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tb.text_frame.vertical_anchor = kw.pop("anchor", MSO_ANCHOR.TOP)
    nowrap = kw.pop("nowrap", False)
    fill_tf(tb.text_frame, lines, base, color, lh, **kw)
    if nowrap:
        tb.text_frame.word_wrap = False
    return tb


def rect(sl, x, y, w, h, fill=None, line=None, radius=0.0, lw=1.0):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = sl.shapes.add_shape(shape, IN(x), IN(y), IN(w), IN(h))
    if radius:
        sh.adjustments[0] = min(0.5, radius / max(1.0, min(w, h)))
    sh.shadow.inherit = False
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor.from_string(fill)
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = RGBColor.from_string(line)
        sh.line.width = Pt(lw / 2)
    else:
        sh.line.fill.background()
    sh.text_frame.word_wrap = True
    return sh


def png_size(p: Path) -> tuple[int, int]:
    return struct.unpack(">II", p.read_bytes()[16:24])


# ── 표 (네이티브 pptx 표. 기본 파란 스타일을 꺼야 어두운 덱에 얹힌다)
NO_STYLE = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"


def _plain_table(tbl):
    pr = tbl._tbl.tblPr
    pr.set("firstRow", "0")
    pr.set("bandRow", "0")
    for el in pr.findall(qn("a:tableStyleId")):
        pr.remove(el)
    sid = pr.makeelement(qn("a:tableStyleId"), {})
    sid.text = NO_STYLE
    pr.append(sid)


def _cell_bottom(cell, color: str, w_pt: float = 0.5):
    tcPr = cell._tc.get_or_add_tcPr()
    ln = tcPr.makeelement(qn("a:lnB"), {"w": str(int(w_pt * 12700)), "cap": "flat", "algn": "ctr"})
    sf = ln.makeelement(qn("a:solidFill"), {})
    sf.append(ln.makeelement(qn("a:srgbClr"), {"val": color}))
    ln.append(sf)
    tcPr.insert(0, ln)


# ────────────────────────────────────────────────────── 블록 (측정 + 그리기)

@dataclass
class Block:
    h: float
    draw: Callable[[object, float, float], None]
    mt: float = 0.0
    mb: float = 0.0
    bottom: bool = False        # margin-top:auto — 콘텐츠 영역 바닥에 붙인다


def para(n: Node, w: float, size: float, color: str, lh: float, *,
         mt=0.0, mb=0.0, bold=False, align=PP_ALIGN.LEFT, tracking=0.0) -> Block:
    st = inline_style(n.attrs.get("style", ""))
    size = st.get("size", size)
    color = st.get("color", color)
    mt = css_len_of(n, "margin-top", mt)
    mb = css_len_of(n, "margin-bottom", mb)
    lines = lines_of(n, {"size": size, "color": color, "bold": bold})
    h = text_h(lines, w, size, lh)
    return Block(h, lambda sl, x, y: text_box(sl, x, y, w, h + size, lines, size, color, lh,
                                              bold=bold, align=align, tracking=tracking),
                 mt, mb)


def css_len_of(n: Node, prop: str, default: float) -> float:
    m = re.search(rf"{prop}\s*:\s*([\d.]+)vmin", n.attrs.get("style", ""))
    if m:
        return float(m.group(1)) * VMIN
    if re.search(rf"{prop}\s*:\s*auto", n.attrs.get("style", "")):
        return -1.0
    return default


# ── .card
def card_block(n: Node, w: float) -> Block:
    pad = 2.7 * VMIN
    gap = 1.1 * VMIN
    inner = w - 2 * pad
    on = "on" in n.cls
    parts: list[tuple[list[Line], float, str, float, bool]] = []
    for k in n.kids():
        if "t" in k.cls:
            st = inline_style(k.attrs.get("style", ""))
            sz = st.get("size", 2.1 * VMIN)
            parts.append((lines_of(k, {"size": sz, "color": st.get("color", C["text"]), "bold": True}),
                          sz, C["text"], 1.25, True))
        elif "d" in k.cls:
            st = inline_style(k.attrs.get("style", ""))
            sz = st.get("size", 1.8 * VMIN)
            col = st.get("color", C["muted"])
            lh = 1.5 if sz > 2.2 * VMIN else 1.6
            parts.append((lines_of(k, {"size": sz, "color": col}), sz, col, lh, False))
    h = 2 * pad + sum(text_h(p[0], inner, p[1], p[3]) for p in parts) + gap * (len(parts) - 1)

    def draw(sl, x, y):
        rect(sl, x, y, w, h, fill=C["on"] if on else C["surface"],
             line=C["kb"] if on else C["line"], radius=1.4 * VMIN)
        cy = y + pad
        for lines, sz, col, lh, bold in parts:
            ph = text_h(lines, inner, sz, lh)
            text_box(sl, x + pad, cy, inner, ph + sz, lines, sz, col, lh, bold=bold)
            cy += ph + gap
    return Block(h, draw)


def cards_block(n: Node, w: float) -> Block:
    cards = [k for k in n.kids() if "card" in k.cls]
    cols = 2 if "two" in n.cls else 3
    gap = 2 * VMIN
    cw = (w - gap * (cols - 1)) / cols
    blocks = [card_block(c, cw) for c in cards]
    h = max(b.h for b in blocks)

    def draw(sl, x, y):
        for i, b in enumerate(blocks):
            bx = x + (i % cols) * (cw + gap)
            by = y + (i // cols) * (h + gap)
            rect(sl, bx, by, cw, h, fill=C["on"] if "on" in cards[i].cls else C["surface"],
                 line=C["kb"] if "on" in cards[i].cls else C["line"], radius=1.4 * VMIN)
            _card_text(sl, cards[i], bx, by, cw)
    rows = math.ceil(len(blocks) / cols)
    return Block(h * rows + gap * (rows - 1), draw, mt=css_len_of(n, "margin-top", 1 * VMIN))


def _card_text(sl, n: Node, x: float, y: float, w: float):
    pad, gap, inner = 2.7 * VMIN, 1.1 * VMIN, w - 2 * 2.7 * VMIN
    cy = y + pad
    for k in n.kids():
        st = inline_style(k.attrs.get("style", ""))
        if "t" in k.cls:
            sz = st.get("size", 2.1 * VMIN)
            lines = lines_of(k, {"size": sz, "color": st.get("color", C["text"]), "bold": True})
            lh, col, bold = 1.25, st.get("color", C["text"]), True
        elif "d" in k.cls:
            sz = st.get("size", 1.8 * VMIN)
            col = st.get("color", C["muted"])
            lines = lines_of(k, {"size": sz, "color": col})
            lh, bold = (1.5 if sz > 2.2 * VMIN else 1.6), False
        else:
            continue
        ph = text_h(lines, inner, sz, lh)
        text_box(sl, x + pad, cy, inner, ph + sz, lines, sz, col, lh, bold=bold)
        cy += ph + gap


# ── table
def table_block(n: Node, w: float) -> Block:
    tight = "tight" in n.cls
    base = 1.7 * VMIN if tight else 1.85 * VMIN
    th_sz = 1.6 * VMIN
    px_, py_ = (0.9 * VMIN, 1.0 * VMIN) if tight else (1.1 * VMIN, 1.3 * VMIN)
    rows = [r for r in n.kids() if r.tag == "tr"] or \
           [r for k in n.kids() for r in k.kids() if r.tag == "tr"]

    ncol = max(len(r.kids()) for r in rows)
    widths = [None] * ncol
    for r in rows:
        for j, c in enumerate(r.kids()):
            if m := re.search(r"width\s*:\s*([\d.]+)%", c.attrs.get("style", "")):
                widths[j] = w * float(m.group(1)) / 100
    rest = [i for i, x in enumerate(widths) if x is None]
    left = w - sum(x for x in widths if x)
    for i in rest:
        widths[i] = left / len(rest)

    cells: list[list[tuple[list[Line], float, str, int, bool]]] = []
    heights: list[float] = []
    for r in rows:
        key = "key" in r.cls
        row: list = []
        for j, c in enumerate(r.kids()):
            head = c.tag == "th"
            st = inline_style(c.attrs.get("style", ""))
            sz = st.get("size", th_sz if head else base)
            col = st.get("color", C["dim"] if head else (C["kb"] if key else C["text"]))
            if "blank" in c.cls:
                col = C["line"]
            al = PP_ALIGN.LEFT
            if "n" in c.cls or "text-align:right" in c.attrs.get("style", "").replace(" ", ""):
                al = PP_ALIGN.RIGHT
            if "text-align:center" in c.attrs.get("style", "").replace(" ", ""):
                al = PP_ALIGN.CENTER
            lines = lines_of(c, {"size": sz, "color": col, "bold": key or head})
            row.append((lines, sz, col, al, key or head))
        cells.append(row)
        heights.append(max(
            text_h(lines, widths[j] - 2 * px_, sz, 1.45) for j, (lines, sz, *_ ) in enumerate(row)
        ) + 2 * py_)
    total = sum(heights)

    def draw(sl, x, y):
        gf = sl.shapes.add_table(len(rows), ncol, IN(x), IN(y), IN(w), IN(total))
        tbl = gf.table
        _plain_table(tbl)
        for j in range(ncol):
            tbl.columns[j].width = IN(widths[j])
        for i, row in enumerate(cells):
            tbl.rows[i].height = IN(heights[i])
            for j in range(ncol):
                cell = tbl.cell(i, j)
                cell.fill.background()
                cell.margin_left = cell.margin_right = IN(px_)
                cell.margin_top = cell.margin_bottom = IN(py_)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                _cell_bottom(cell, C["line"])
                if j < len(row):
                    lines, sz, col, al, bold = row[j]
                    fill_tf(cell.text_frame, lines, sz, col, 1.45, align=al, bold=bold)
    return Block(total, draw, mt=css_len_of(n, "margin-top", 1 * VMIN))


# ── .formula
def formula_block(n: Node, w: float) -> Block:
    st = inline_style(n.attrs.get("style", ""))
    size = st.get("size", 2.0 * VMIN)
    px_, py_ = 2.4 * VMIN, 2.2 * VMIN
    lines = lines_of(n, {"size": size, "color": C["text"], "mono": True})
    inner = w - 2 * px_ - 3
    th = text_h(lines, inner, size, 1.85)
    h = th + 2 * py_

    def draw(sl, x, y):
        rect(sl, x, y, w, h, fill=C["surface"], line=C["line"], radius=0.9 * VMIN)
        rect(sl, x, y, 4, h, fill=C["kb"])
        text_box(sl, x + px_, y + py_, inner, th + size, lines, size, C["text"], 1.85, mono=True)
    return Block(h, draw,
                 mt=css_len_of(n, "margin-top", 1.8 * VMIN),
                 mb=css_len_of(n, "margin-bottom", 1.8 * VMIN))


# ── .stack (L0~L4)
def stack_block(n: Node, w: float) -> Block:
    gap = 1.1 * VMIN
    px_, py_ = 2.1 * VMIN, 1.6 * VMIN
    idw, icol = 7 * VMIN, 2 * VMIN
    inner = w - 2 * px_ - idw - icol - 3
    rows = []
    for lay in n.kids():
        opt = "opt" in lay.cls
        d = next(k for k in lay.kids() if "d" in k.cls)
        i = next(k for k in lay.kids() if "id" in k.cls)
        lines = lines_of(d, {"size": 1.8 * VMIN, "color": C["dim"] if opt else C["muted"]})
        if opt:   # HTML은 레이어 전체에 opacity .72 — 강조 런까지 같이 죽인다
            lines = [[(t, {**st, "color": C["muted"]} if st.get("color") == C["text"] else st)
                      for t, st in ln] for ln in lines]
        h = max(text_h(lines, inner, 1.8 * VMIN, 1.5) + 2 * py_, 2.1 * VMIN * 1.5 + 2 * py_)
        rows.append((txt(i), lines, h, opt))
    total = sum(r[2] for r in rows) + gap * (len(rows) - 1)

    def draw(sl, x, y):
        cy = y
        for label, lines, h, opt in rows:
            rect(sl, x, cy, w, h, fill=C["surface"], line=C["line"], radius=0.9 * VMIN)
            rect(sl, x, cy, 4, h, fill=C["dim"] if opt else C["kb"])
            text_box(sl, x + px_ + 3, cy, idw, h, [[(label, {})]], 2.1 * VMIN,
                     C["dim"] if opt else C["kb"], 1.2, bold=True, anchor=MSO_ANCHOR.MIDDLE)
            text_box(sl, x + px_ + 3 + idw + icol, cy, inner, h, lines, 1.8 * VMIN,
                     C["dim"] if opt else C["muted"], 1.5, anchor=MSO_ANCHOR.MIDDLE)
            cy += h + gap
    return Block(total, draw, mt=css_len_of(n, "margin-top", 1 * VMIN))


# ── .flow (파이프라인 4단계)
def flow_block(n: Node, w: float) -> Block:
    steps = [k for k in n.kids() if "step" in k.cls]
    gap, chev = 1.4 * VMIN, 2.6 * VMIN
    sw = (w - gap * 2 * (len(steps) - 1) - chev * (len(steps) - 1)) / len(steps)
    px_, py_, ig = 1.7 * VMIN, 1.9 * VMIN, 0.8 * VMIN
    inner = sw - 2 * px_
    body = []
    for s in steps:
        n_, t_, d_ = (next(k for k in s.kids() if c in k.cls) for c in ("n", "t", "d"))
        dl = lines_of(d_, {"size": 1.6 * VMIN, "color": C["muted"]})
        h = (2 * py_ + 1.5 * VMIN * 1.2 + 1.9 * VMIN * 1.25
             + text_h(dl, inner, 1.6 * VMIN, 1.5) + 2 * ig)
        body.append((txt(n_), lines_of(t_, {"size": 1.9 * VMIN, "color": C["text"], "bold": True}),
                     dl, h, "warn" in s.cls))
    total = max(b[3] for b in body)

    def draw(sl, x, y):
        cx = x
        for i, (num, tl, dl, _h, warn) in enumerate(body):
            rect(sl, cx, y, sw, total, fill=C["surface"],
                 line=C["bad"] if warn else C["line"], radius=1.0 * VMIN)
            cy = y + py_
            text_box(sl, cx + px_, cy, inner, 1.5 * VMIN * 1.6,
                     [[(num, {})]], 1.5 * VMIN, C["bad"] if warn else C["kb"], 1.2,
                     bold=True, tracking=0.1)
            cy += 1.5 * VMIN * 1.2 + ig
            th = text_h(tl, inner, 1.9 * VMIN, 1.25)
            text_box(sl, cx + px_, cy, inner, th + 1.9 * VMIN, tl, 1.9 * VMIN, C["text"], 1.25, bold=True)
            cy += th + ig
            dh = text_h(dl, inner, 1.6 * VMIN, 1.5)
            text_box(sl, cx + px_, cy, inner, dh + 1.6 * VMIN, dl, 1.6 * VMIN, C["muted"], 1.5)
            cx += sw
            if i < len(body) - 1:
                text_box(sl, cx + gap, y, chev, total, [[("›", {})]], 2.2 * VMIN, C["dim"], 1.0,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                cx += gap * 2 + chev
    return Block(total, draw,
                 mt=css_len_of(n, "margin-top", 2 * VMIN), mb=css_len_of(n, "margin-bottom", 2 * VMIN))


# ── dl.kv / dl.meta
def dl_block(n: Node, w: float, *, meta=False) -> Block:
    rgap, cgap = (0.9 * VMIN, 2.4 * VMIN) if meta else (1.2 * VMIN, 2.2 * VMIN)
    dt_sz = 1.65 * VMIN if meta else 1.85 * VMIN
    dd_sz = 1.9 * VMIN if meta else 1.85 * VMIN
    dt_col = C["dim"] if meta else C["kb"]
    dd_col = C["text"] if meta else C["muted"]
    lh = 1.3 if meta else 1.55
    pairs, dtw = [], 0.0
    kids = n.kids()
    for i in range(0, len(kids) - 1, 2):
        dt, dd = kids[i], kids[i + 1]
        label = txt(dt)
        dtw = max(dtw, text_w(label, dt_sz))
        pairs.append((label, dd))
    dtw = 11 * VMIN if meta else dtw + 10
    ddw = w - dtw - cgap
    rows = []
    for label, dd in pairs:
        blanks = dd.find("blankfill")
        lines = lines_of(dd, {"size": dd_sz, "color": dd_col}, skip=lambda k: "blankfill" in k.cls)
        h = max(text_h(lines, ddw, dd_sz, lh), dd_sz * lh)
        rows.append((label, lines, h, blanks))
    total = sum(r[2] for r in rows) + rgap * (len(rows) - 1)

    def draw(sl, x, y):
        cy = y
        for label, lines, h, blanks in rows:
            # .kv dt는 white-space:nowrap — 폭 추정이 짧으면 '판단값 표/시'처럼 갈라진다
            text_box(sl, x, cy, dtw, h, [[(label, {})]], dt_sz, dt_col, lh,
                     bold=not meta, nowrap=not meta)
            if blanks:
                bx = x + dtw + cgap
                for _ in blanks:
                    rect(sl, bx, cy + dd_sz * lh * 0.92, 24 * VMIN, 1.2, fill=C["dim"])
                    bx += 24 * VMIN + 1.2 * VMIN
            if any(t.strip() for ln in lines for t, _ in ln):
                text_box(sl, x + dtw + cgap, cy, ddw, h + dd_sz, lines, dd_sz, dd_col, lh)
            cy += h + rgap
    return Block(total, draw, mt=css_len_of(n, "margin-top", (3.4 if meta else 1.4) * VMIN))


# ── .tenure (전세/월세 점유형태 막대)
def tenure_block(n: Node, w: float) -> Block:
    gap = 3 * VMIN
    tens = n.kids()
    flex = [2.0 if "flex:2" in t.attrs.get("style", "").replace(" ", "") else 1.0 for t in tens]
    unit = (w - gap * (len(tens) - 1)) / sum(flex)
    bars_h, lab_h = 24 * VMIN, 1.9 * VMIN * 1.3 + 1.2 * VMIN
    total = bars_h + lab_h

    def draw(sl, x, y):
        cx = x
        for t, fl in zip(tens, flex):
            cw = unit * fl
            group = t.find("bars")
            if group:
                bars = group[0].kids()
                bw, bg = 4.4 * VMIN, 1.4 * VMIN
                bx = cx + (cw - (len(bars) * bw + (len(bars) - 1) * bg)) / 2
                for b in bars:
                    pct = float(re.search(r"height\s*:\s*([\d.]+)%", b.attrs["style"]).group(1))
                    bh = bars_h * pct / 100
                    by = y + bars_h - bh
                    rect(sl, bx, by, bw, bh,
                         fill=C["kb"] if "b24" in b.cls else "4A4438", radius=0.5 * VMIN)
                    text_box(sl, bx - bw / 2, by - 2.8 * VMIN, bw * 2, 2.4 * VMIN,
                             [[(txt(b), {})]], 1.7 * VMIN, C["muted"], 1.2, align=PP_ALIGN.CENTER)
                    bx += bw + bg
                lab = t.find("lab")[0]
                text_box(sl, cx, y + bars_h + 1.2 * VMIN, cw, lab_h, [[(txt(lab), {})]],
                         1.9 * VMIN, C["muted"], 1.3, align=PP_ALIGN.CENTER)
            else:
                p = t.find("lead")[0]
                sz = inline_style(p.attrs.get("style", "")).get("size", 2.2 * VMIN)
                lines = lines_of(p, {"size": sz, "color": C["muted"]})
                th = text_h(lines, cw, sz, 1.65)
                text_box(sl, cx, y + (total - th) / 2, cw, th + sz, lines, sz, C["muted"], 1.65)
            cx += cw + gap
    return Block(total, draw, mt=3 * VMIN, mb=3 * VMIN)


# ── .flip (전/후 큰 숫자)
def flip_block(n: Node, w: float) -> Block:
    gap = 3 * VMIN
    arrow = 4 * VMIN
    sw = (w - 2 * gap - arrow) / 2
    pad = 3.2 * VMIN
    states = [k for k in n.kids() if "state" in k.cls]
    body = []
    for s in states:
        good = "a" in s.cls
        col = C["good"] if good else C["kb"]
        cap = next(k for k in s.kids() if "cap" in k.cls)
        num = next(k for k in s.kids() if "num" in k.cls)
        sub = next(k for k in s.kids() if "sub" in k.cls)
        cl = lines_of(cap, {"size": 1.85 * VMIN, "color": C["muted"]})
        ch_ = text_h(cl, sw - 2 * pad, 1.85 * VMIN, 1.35)
        body.append((good, col, cl, ch_, lines_of(num, {"size": 6.2 * VMIN, "color": col, "bold": True}),
                     lines_of(sub, {"size": 1.95 * VMIN, "color": col, "bold": True})))
    inner = max(b[3] for b in body) + 1.3 * VMIN + 6.2 * VMIN + 1.3 * VMIN + 1.95 * VMIN * 1.3
    total = inner + 2 * pad

    def draw(sl, x, y):
        for i, (good, col, cl, ch_, nl, sl_) in enumerate(body):
            bx = x + i * (sw + gap + arrow + gap)
            rect(sl, bx, y, sw, total, fill=C["surface"],
                 line=C["line"] if good else C["kb"], radius=1.4 * VMIN)
            cy = y + pad
            text_box(sl, bx + pad, cy, sw - 2 * pad, ch_ + 1.85 * VMIN, cl, 1.85 * VMIN,
                     C["muted"], 1.35, align=PP_ALIGN.CENTER)
            cy += ch_ + 1.3 * VMIN
            text_box(sl, bx + pad, cy, sw - 2 * pad, 6.2 * VMIN * 1.15, nl, 6.2 * VMIN, col, 1.0,
                     bold=True, align=PP_ALIGN.CENTER)
            cy += 6.2 * VMIN + 1.3 * VMIN
            text_box(sl, bx + pad, cy, sw - 2 * pad, 1.95 * VMIN * 1.6, sl_, 1.95 * VMIN, col, 1.3,
                     bold=True, align=PP_ALIGN.CENTER)
        text_box(sl, x + sw + gap, y, arrow, total, [[("→", {})]], 4 * VMIN, C["dim"], 1.0,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return Block(total, draw, mt=3 * VMIN, mb=3 * VMIN)


# ── .div-chart (0을 기준으로 좌우로 뻗는 막대)
def divchart_block(n: Node, w: float) -> Block:
    labw, gap = 15 * VMIN, 2 * VMIN
    trw = w - labw - gap
    rh, rgap = 5 * VMIN, 2 * VMIN
    rows = n.kids()
    total = len(rows) * rh + (len(rows) - 1) * rgap

    def draw(sl, x, y):
        cy = y
        for i, r in enumerate(rows):
            rl = r.find("rl")[0]
            lines = lines_of(rl, {"size": 1.85 * VMIN, "color": C["muted"]})
            lh_ = text_h(lines, labw, 1.85 * VMIN, 1.35)
            text_box(sl, x, cy + (rh - lh_) / 2, labw, lh_ + 2 * VMIN, lines, 1.85 * VMIN,
                     C["muted"], 1.35, align=PP_ALIGN.RIGHT)
            tx = x + labw + gap
            rect(sl, tx, cy, trw, rh, fill=C["surface"], line=C["line"], radius=0.6 * VMIN)
            zx = tx + trw * 0.12
            rect(sl, zx, cy - 0.6 * VMIN, 2, rh + 1.2 * VMIN, fill=C["dim"])
            if i == 0:
                text_box(sl, zx - 3 * VMIN, cy - 3.4 * VMIN, 6 * VMIN, 2.2 * VMIN,
                         [[("0", {})]], 1.5 * VMIN, C["dim"], 1.2, align=PP_ALIGN.CENTER)
            f = r.find("fill")[0]
            pct = float(re.search(r"width\s*:\s*([\d.]+)%", f.attrs["style"]).group(1))
            bw = max(trw * pct / 100, 3)
            fy, fh = cy + 0.7 * VMIN, rh - 1.4 * VMIN
            v = f.find("v")[0]
            neg, big = "neg" in f.cls, "big" in f.cls
            if neg:
                rect(sl, zx - bw, fy, bw, fh, fill=C["good"], radius=0.4 * VMIN)
                text_box(sl, zx - bw - 1.2 * VMIN - 12 * VMIN, cy, 12 * VMIN, rh,
                         [[(txt(v), {})]], 1.95 * VMIN, C["good"], 1.2, bold=True,
                         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
            else:
                rect(sl, zx, fy, bw, fh, fill=C["kb"] if big else C["kb_deep"], radius=0.4 * VMIN)
                if big:
                    text_box(sl, zx, cy, bw - 1.4 * VMIN, rh, [[(txt(v), {})]], 1.95 * VMIN,
                             C["kb_ink"], 1.2, bold=True, align=PP_ALIGN.RIGHT,
                             anchor=MSO_ANCHOR.MIDDLE)
                else:
                    text_box(sl, zx + bw + 1.2 * VMIN, cy, 12 * VMIN, rh, [[(txt(v), {})]],
                             1.95 * VMIN, C["text"], 1.2, bold=True, anchor=MSO_ANCHOR.MIDDLE)
            cy += rh + rgap
    return Block(total, draw, mt=2 * VMIN + 1.4 * VMIN, mb=1 * VMIN)


# ── .shot (스크린샷을 흰 카드에 얹는다)
def shot_block(n: Node, w: float) -> Block:
    pad = 2 * VMIN if "pad2" in n.cls else 1.4 * VMIN
    img = n.find("img")[0]
    src = (DECK.parent / img.attrs["src"]).resolve()
    iw, ih = png_size(src)
    dw = w - 2 * pad
    dh = dw * ih / iw
    h = dh + 2 * pad

    def draw(sl, x, y):
        rect(sl, x, y, w, h, fill="FFFFFF", radius=1.2 * VMIN)
        sl.shapes.add_picture(str(src), IN(x + pad), IN(y + pad), IN(dw), IN(dh))
    return Block(h, draw)


# ── .split (2단)
def split_block(n: Node, w: float) -> Block:
    gap = 3 * VMIN
    ratio = (0.85, 1.15) if "wide-r" in n.cls else (1.0, 1.0)
    cols = n.kids()
    ws = [(w - gap) * r / sum(ratio) for r in ratio]
    # 단(段)이 곧 하나의 블록일 수 있다(표·카드·dl). 이걸 컨테이너로만 보면 통째로 사라진다.
    stacks = [([b] if (b := block_of(c, ws[i])) else stack_of(c, ws[i])) for i, c in enumerate(cols)]
    hs = [sum(b.h for b in s) + sum(max(0, b.mt) + max(0, b.mb) for b in s[1:]) for s in stacks]
    total = max(hs)

    def draw(sl, x, y):
        cx = x
        for i, s in enumerate(stacks):
            place(sl, s, cx, y + (total - hs[i]) / 2)
            cx += ws[i] + gap
    return Block(total, draw, mt=css_len_of(n, "margin-top", 1 * VMIN))


# ────────────────────────────────────────────────────── 블록 디스패치

def block_of(n: Node, w: float) -> Block | None:
    c = n.cls
    if "kicker" in c:
        return para(n, w, 1.65 * VMIN, C["kb"], 1.3, mb=1.5 * VMIN, bold=True, tracking=0.16)
    if n.tag == "h1":
        return para(n, w, 7.6 * VMIN, C["text"], 1.15, bold=True, tracking=-0.02)
    if n.tag == "h2":
        return para(n, w, 4.2 * VMIN, C["text"], 1.25, mb=2.4 * VMIN, bold=True, tracking=-0.02)
    if "lead" in c:
        return para(n, w, 2.4 * VMIN, C["muted"], 1.65)
    if "note" in c:
        return para(n, w, 1.65 * VMIN, C["dim"], 1.6, mt=2.2 * VMIN)
    if "shotcap" in c:
        return para(n, w, 1.5 * VMIN, C["dim"], 1.5, mt=1.1 * VMIN)
    if "mark" in c:
        return para(n, w, 2.1 * VMIN, C["muted"], 1.3, mb=2.6 * VMIN, tracking=0.06)
    if "foot" in c:
        b = para(n, w, 1.75 * VMIN, C["dim"], 1.5, mt=css_len_of(n, "margin-top", -1.0))
        return Block(b.h, b.draw, mt=max(0.0, b.mt), bottom=b.mt < 0)
    if "q" in c:
        return quote_block(n, w)
    if "meta" in c:
        return dl_block(n, w, meta=True)
    if "kv" in c:
        return dl_block(n, w)
    if "cards" in c:
        return cards_block(n, w)
    if "card" in c:
        b = card_block(n, w)
        return Block(b.h, lambda sl, x, y: (rect(sl, x, y, w, b.h,
                                                 fill=C["on"] if "on" in c else C["surface"],
                                                 line=C["kb"] if "on" in c else C["line"],
                                                 radius=1.4 * VMIN), _card_text(sl, n, x, y, w))[0],
                     mt=css_len_of(n, "margin-top", 0.0))
    if n.tag == "table":
        return table_block(n, w)
    if "formula" in c:
        return formula_block(n, w)
    if "stack" in c:
        return stack_block(n, w)
    if "flow" in c:
        return flow_block(n, w)
    if "tenure" in c:
        return tenure_block(n, w)
    if "flip" in c:
        return flip_block(n, w)
    if "div-chart" in c:
        return divchart_block(n, w)
    if "shot" in c:
        return shot_block(n, w)
    if "split" in c:
        return split_block(n, w)
    if n.tag == "div" and n.kids():          # .split 안의 무명 래퍼
        blocks = stack_of(n, w)
        h = sum(b.h for b in blocks) + sum(max(0, b.mt) + max(0, b.mb) for b in blocks[1:])
        return Block(h, lambda sl, x, y: place(sl, blocks, x, y))
    if n.tag == "p":
        return para(n, w, 2.4 * VMIN, C["muted"], 1.65)
    return None


def quote_block(n: Node, w: float) -> Block:
    padl = 2.6 * VMIN
    lines = lines_of(n, {"size": 2.9 * VMIN, "color": C["text"], "bold": True})
    h = text_h(lines, w - padl, 2.9 * VMIN, 1.5)

    def draw(sl, x, y):
        rect(sl, x, y, 5, h, fill=C["kb"])
        text_box(sl, x + padl, y, w - padl, h + 2.9 * VMIN, lines, 2.9 * VMIN, C["text"], 1.5, bold=True)
    return Block(h, draw, mt=3.4 * VMIN)


def stack_of(parent: Node, w: float) -> list[Block]:
    out = []
    for k in parent.kids():
        if b := block_of(k, w):
            out.append(b)
    return out


def place(sl, blocks: list[Block], x: float, y: float):
    for i, b in enumerate(blocks):
        if i:
            y += max(0.0, b.mt)
        b.draw(sl, x, y)
        y += b.h + max(0.0, b.mb)


# ────────────────────────────────────────────────────────────── 빌드

def build(html: Path = DECK, out: Path = OUT) -> int:
    p = _Parser()
    p.feed(html.read_text(encoding="utf-8"))
    slides = [n for n in p.root.find("slide")]
    if not slides:
        sys.exit("슬라이드를 찾지 못했습니다 — 덱 구조가 바뀌었는지 확인하세요.")

    prs = Presentation()
    prs.slide_width, prs.slide_height = IN(W), IN(H)
    blank = prs.slide_layouts[6]

    overflow = []
    for i, sec in enumerate(slides, 1):
        sl = prs.slides.add_slide(blank)
        sl.background.fill.solid()
        sl.background.fill.fore_color.rgb = RGBColor.from_string(C["bg"])

        anim = next((k for k in sec.kids() if "anim" in k.cls), sec)
        blocks = stack_of(anim, CW)
        flow = [b for b in blocks if not b.bottom]
        pinned = [b for b in blocks if b.bottom]

        used = sum(b.h for b in flow) + sum(max(0, b.mt) + max(0, b.mb) for b in flow[1:])
        if used > CH:
            overflow.append((i, round(used - CH)))
        # 표지(.foot이 margin-top:auto)는 위에서 시작하고 푸터만 바닥에 붙는다 — 브라우저와 같다.
        y0 = PAD_Y if pinned else PAD_Y + max(0.0, (CH - used) / 2)
        place(sl, flow, PAD_X, y0)
        for b in pinned:
            b.draw(sl, PAD_X, PAD_Y + CH - b.h)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    for i, over in overflow:
        print(f"  ⚠ 슬라이드 {i}: 추정 높이가 {over}px 넘칩니다 — 확인하세요.")
    print(f"{out.relative_to(ROOT)}  {out.stat().st_size // 1024} KB  ·  슬라이드 {len(slides)}장 (편집 가능)")
    return len(slides)


if __name__ == "__main__":
    build()
