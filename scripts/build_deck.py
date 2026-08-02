#!/usr/bin/env python3
"""기술설명서 덱(HTML) → 슬라이드 PNG → .pptx + .pdf.

    uv pip install -p .venv -e ".[docs]"      # python-pptx (문서 빌드 전용)
    .venv/bin/python scripts/build_deck.py

`docs/deck/*.html`이 원본이다. PNG·PPTX·PDF는 전부 산출물이라 `dist-docs/`로 나가고
저장소가 추적하지 않는다 — 원본만 고치고 이 스크립트를 다시 돌린다.

왜 이미지 경유인가: 덱은 vmin 기반 CSS라 브라우저 렌더가 유일한 정본이다. pptx로
텍스트를 옮기면 조판이 달라져 '보이는 것'과 '제출한 것'이 갈라진다. PPTX와 PDF가
**같은 PNG**를 쓰므로 두 파일이 어긋날 수 없다.
"""
from __future__ import annotations

import re
import shutil
import subprocess
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
DECK = ROOT / "docs" / "deck" / "온전_기술설명서.html"
OUT = ROOT / "dist-docs"
PNG_DIR = OUT / "slides"
W, H = 1920, 1080  # 16:9. vmin = 1080 → 화면에서 보던 배치가 그대로 나온다.

CHROME_CANDIDATES = [
    "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
    "/Applications/Chromium.app/Contents/MacOS/Chromium",
    "/Applications/Microsoft Edge.app/Contents/MacOS/Microsoft Edge",
]


def find_chrome() -> str:
    for p in CHROME_CANDIDATES:
        if Path(p).exists():
            return p
    for name in ("google-chrome", "chromium", "chromium-browser"):
        if found := shutil.which(name):
            return found

    sys.exit("Chrome/Chromium을 찾지 못했습니다 — CHROME_CANDIDATES에 경로를 추가하세요.")


def slide_count(html: str) -> int:
    # <section class="slide ..."> 개수. 덱의 슬라이드 정의는 이 한 패턴뿐이다.
    return len(re.findall(r'<section class="slide', html))


def chrome_run(chrome: str, *args: str) -> None:
    # --force-prefers-reduced-motion: 덱의 등장 애니메이션을 끈다. 안 끄면 캡처 시점에
    # 따라 요소가 opacity 0에서 잡혀 슬라이드가 반쯤 빈 채로 나온다.
    base = [
        chrome, "--headless=new", "--disable-gpu", "--hide-scrollbars",
        "--force-prefers-reduced-motion", "--allow-file-access-from-files",
        "--virtual-time-budget=4000", "--no-first-run", "--no-default-browser-check",
    ]
    subprocess.run([*base, *args], check=True, capture_output=True)


def main() -> None:
    if not DECK.exists():
        sys.exit(f"덱 원본이 없습니다: {DECK}")
    chrome = find_chrome()
    n = slide_count(DECK.read_text(encoding="utf-8"))
    if n == 0:
        sys.exit("슬라이드를 하나도 찾지 못했습니다 — 덱 구조가 바뀌었는지 확인하세요.")

    PNG_DIR.mkdir(parents=True, exist_ok=True)
    pngs: list[Path] = []
    for i in range(1, n + 1):
        png = PNG_DIR / f"slide-{i:02d}.png"
        chrome_run(chrome, f"--window-size={W},{H}", f"--screenshot={png}",
                   f"{DECK.as_uri()}?s={i}")
        if not png.exists() or png.stat().st_size < 5_000:
            sys.exit(f"슬라이드 {i} 캡처가 비었습니다({png}) — 렌더 실패를 조용히 넘기지 않습니다.")
        pngs.append(png)
        print(f"  slide {i:02d}/{n}  {png.stat().st_size // 1024:>4} KB")

    # ── PPTX: 16:9 슬라이드에 PNG를 전면 배치
    from pptx import Presentation
    from pptx.util import Emu, Inches

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]  # 빈 레이아웃 — 플레이스홀더가 이미지 위에 얹히지 않는다
    for png in pngs:
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(str(png), Emu(0), Emu(0),
                             width=prs.slide_width, height=prs.slide_height)
    pptx_path = OUT / "온전_기술설명서.pptx"
    prs.save(pptx_path)

    # ── PDF: 같은 PNG를 한 장/페이지로 깔고 Chrome으로 인쇄
    imgs = "\n".join(f'<img src="{p.as_uri()}">' for p in pngs)
    printable = PNG_DIR / "_print.html"
    printable.write_text(
        "<!DOCTYPE html><meta charset='utf-8'>"
        "<style>@page{size:13.333in 7.5in;margin:0}"
        "html,body{margin:0;padding:0;background:#1b1917}"
        "img{display:block;width:13.333in;height:7.5in;page-break-after:always}"
        "img:last-child{page-break-after:auto}</style>" + imgs,
        encoding="utf-8",
    )
    pdf_path = OUT / "온전_기술설명서.pdf"
    chrome_run(chrome, "--no-pdf-header-footer", f"--print-to-pdf={pdf_path}",
               printable.as_uri())

    for path in (pptx_path, pdf_path):
        print(f"{path.relative_to(ROOT)}  {path.stat().st_size // 1024} KB")
    print(f"슬라이드 {n}장 · PNG {PNG_DIR.relative_to(ROOT)}/")


if __name__ == "__main__":
    main()
